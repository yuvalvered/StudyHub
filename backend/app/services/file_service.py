"""
File service for processing uploaded files (PDF, Word, PowerPoint, Excel text extraction)
"""
from pathlib import Path
from typing import Optional
import logging
import re

try:
    import fitz  # pymupdf
    PYMUPDF_AVAILABLE = True
except ImportError:
    PYMUPDF_AVAILABLE = False
    logger_temp = logging.getLogger(__name__)
    logger_temp.warning("pymupdf not installed, falling back to pdfplumber.")

try:
    import pytesseract
    import os
    import shutil
    # Use env var if set, otherwise try common Windows paths, then rely on PATH
    _tesseract_cmd = os.environ.get('TESSERACT_CMD', '')
    if not _tesseract_cmd:
        _candidates = [
            r'C:\Program Files\Tesseract-OCR\tesseract.exe',
            r'C:\Program Files (x86)\Tesseract-OCR\tesseract.exe',
            r'C:\Tesseract-OCR\tesseract.exe',
        ]
        for _c in _candidates:
            if os.path.exists(_c):
                _tesseract_cmd = _c
                break
    if _tesseract_cmd:
        pytesseract.pytesseract.tesseract_cmd = _tesseract_cmd
    # Verify tesseract is actually callable
    pytesseract.get_tesseract_version()
    TESSERACT_AVAILABLE = True
except Exception:
    TESSERACT_AVAILABLE = False

try:
    import pdfplumber
    PDFPLUMBER_AVAILABLE = True
except ImportError:
    PDFPLUMBER_AVAILABLE = False

logger = logging.getLogger(__name__)

# Optional imports for Office documents
try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    logger.warning("python-docx not installed. DOCX text extraction disabled.")

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logger.warning("python-pptx not installed. PPTX text extraction disabled.")

try:
    import openpyxl
    XLSX_AVAILABLE = True
except ImportError:
    XLSX_AVAILABLE = False
    logger.warning("openpyxl not installed. XLSX text extraction disabled.")


class FileService:
    """Service for processing file content."""

    @staticmethod
    def fix_hebrew_direction(text: str) -> str:
        """
        Fix Hebrew text direction issues that occur during PDF extraction.

        PDFs with RTL text extract with fully reversed text - both character order
        within words AND word order in sentences. This function fixes both issues.

        Example:
            Input:  "ת וכרעמב תוטלחה" (reversed from PDF)
            Output: "החלטות במערכות" (correct Hebrew)

        Args:
            text: Extracted text that may have reversed Hebrew

        Returns:
            Text with corrected Hebrew direction
        """
        if not text:
            return text

        # Process line by line to handle mixed Hebrew/English content
        lines = text.split('\n')
        fixed_lines = []

        for line in lines:
            # Check if line contains Hebrew
            has_hebrew = any('\u0590' <= c <= '\u05FF' for c in line)

            if has_hebrew:
                # Split into tokens, keeping separators
                # This captures words, spaces, and punctuation separately
                tokens = re.split(r'(\s+)', line)

                # Separate Hebrew and non-Hebrew segments
                # We need to reverse ONLY the Hebrew segments, not English
                segments = []
                current_segment = []
                is_hebrew_segment = None

                for token in tokens:
                    token_has_hebrew = token and any('\u0590' <= c <= '\u05FF' for c in token)

                    # Start new segment if language changes
                    if is_hebrew_segment is None:
                        is_hebrew_segment = token_has_hebrew
                    elif is_hebrew_segment != token_has_hebrew and token.strip():
                        # Language changed - save current segment
                        segments.append((is_hebrew_segment, current_segment))
                        current_segment = []
                        is_hebrew_segment = token_has_hebrew

                    current_segment.append(token)

                # Add last segment
                if current_segment:
                    segments.append((is_hebrew_segment, current_segment))

                # Process each segment
                fixed_tokens = []
                for is_hebrew, segment_tokens in segments:
                    if is_hebrew:
                        # Reverse Hebrew segment order
                        segment_tokens.reverse()

                        # Reverse each Hebrew word's characters
                        for token in segment_tokens:
                            if token and any('\u0590' <= c <= '\u05FF' for c in token):
                                fixed_tokens.append(token[::-1])
                            else:
                                fixed_tokens.append(token)
                    else:
                        # Keep non-Hebrew as is
                        fixed_tokens.extend(segment_tokens)

                # Join tokens
                fixed_line = ''.join(fixed_tokens)

                # Fix cases where single Hebrew letters got separated by spaces
                fixed_line = re.sub(r'([\u0590-\u05FF]{2,}) ([\u0590-\u05FF]{1})\b', r'\1\2', fixed_line)
                fixed_line = re.sub(r'\b([\u0590-\u05FF]{1}) ([\u0590-\u05FF]{2,})', r'\1\2', fixed_line)

                # Clean up multiple spaces
                fixed_line = re.sub(r' +', ' ', fixed_line).strip()
                fixed_lines.append(fixed_line)
            else:
                # Non-Hebrew lines: keep as is
                fixed_lines.append(line)

        return '\n'.join(fixed_lines)

    @staticmethod
    def get_pdf_page_count(file_path: str) -> Optional[int]:
        """Get the number of pages in a PDF file."""
        try:
            if not Path(file_path).exists():
                logger.error(f"PDF file not found: {file_path}")
                return None

            if PYMUPDF_AVAILABLE:
                doc = fitz.open(file_path)
                count = len(doc)
                doc.close()
                return count
            elif PDFPLUMBER_AVAILABLE:
                with pdfplumber.open(file_path) as pdf:
                    return len(pdf.pages)
            return None
        except Exception as e:
            logger.error(f"Error getting page count from PDF {file_path}: {str(e)}")
            return None

    @staticmethod
    def _is_reversed_hebrew(text: str) -> bool:
        """Detect if Hebrew text is stored reversed (characters in wrong order)."""
        # Common Hebrew words that appear reversed in bad PDFs
        reversed_indicators = ['םולש', 'הדות', 'רועיש', 'אשונ', 'רמאמ', 'תואצות', 'אובמ', 'ןכות']
        normal_indicators = ['שלום', 'תודה', 'שיעור', 'נושא', 'מאמר', 'תוצאות', 'מבוא', 'תוכן']
        text_sample = text[:2000]
        reversed_count = sum(1 for w in reversed_indicators if w in text_sample)
        normal_count = sum(1 for w in normal_indicators if w in text_sample)
        return reversed_count > normal_count

    @staticmethod
    def _extract_pdf_with_pymupdf(file_path: str) -> Optional[str]:
        """Extract text from PDF using pymupdf (better Hebrew support)."""
        doc = fitz.open(file_path)
        pages_text = []

        for page in doc:
            blocks = page.get_text("blocks")
            blocks.sort(key=lambda b: (round(b[1] / 10), b[0]))
            page_text = "\n".join(b[4].strip() for b in blocks if b[4].strip())
            if page_text:
                pages_text.append(page_text)

        doc.close()
        full_text = "\n\n".join(pages_text) if pages_text else None

        # Fix direction only if text is actually reversed
        if full_text and FileService._is_reversed_hebrew(full_text):
            full_text = FileService.fix_hebrew_direction(full_text)

        return full_text

    @staticmethod
    def _extract_pdf_with_pdfplumber(file_path: str) -> Optional[str]:
        """Extract text from PDF using pdfplumber (fallback)."""
        text_content = []
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text_content.append(page_text)
        return "\n".join(text_content) if text_content else None

    @staticmethod
    def _extract_pdf_with_ocr(file_path: str) -> Optional[str]:
        """Extract text from scanned PDF using Tesseract OCR (Hebrew + English)."""
        if not PYMUPDF_AVAILABLE or not TESSERACT_AVAILABLE:
            return None

        try:
            from PIL import Image
            import io

            doc = fitz.open(file_path)
            pages_text = []

            for page in doc:
                mat = fitz.Matrix(200 / 72, 200 / 72)
                pix = page.get_pixmap(matrix=mat)
                img = Image.open(io.BytesIO(pix.tobytes("png")))
                # heb+eng - Hebrew right-to-left with English
                page_text = pytesseract.image_to_string(img, lang='heb+eng')
                if page_text.strip():
                    pages_text.append(page_text)

            doc.close()
            return "\n\n".join(pages_text) if pages_text else None

        except Exception as e:
            logger.error(f"OCR failed for {file_path}: {str(e)}")
            return None

    @staticmethod
    def _is_scanned_pdf(text: str, page_count: int) -> bool:
        """Check if PDF is scanned (very little text per page = likely scanned)."""
        if not text or not page_count:
            return True
        chars_per_page = len(text.strip()) / page_count
        return chars_per_page < 100  # less than 100 chars per page = scanned

    @staticmethod
    def extract_pdf_text(file_path: str) -> Optional[str]:
        """Extract text from PDF. Falls back to OCR if scanned."""
        try:
            if not Path(file_path).exists():
                logger.error(f"PDF file not found: {file_path}")
                return None

            logger.info(f"PDF extraction started: pymupdf={PYMUPDF_AVAILABLE}, tesseract={TESSERACT_AVAILABLE}, pdfplumber={PDFPLUMBER_AVAILABLE}")
            full_text = None

            if PYMUPDF_AVAILABLE:
                full_text = FileService._extract_pdf_with_pymupdf(file_path)
                if full_text:
                    logger.info(f"pymupdf extracted {len(full_text)} chars from {file_path}")

            # Fallback to pdfplumber if pymupdf got nothing
            if not full_text and PDFPLUMBER_AVAILABLE:
                full_text = FileService._extract_pdf_with_pdfplumber(file_path)
                if full_text:
                    logger.info(f"pdfplumber extracted {len(full_text)} chars from {file_path}")

            # If result is too short → likely scanned PDF → try OCR
            page_count = FileService.get_pdf_page_count(file_path) or 1
            if FileService._is_scanned_pdf(full_text, page_count) and TESSERACT_AVAILABLE:
                logger.info(f"PDF appears scanned, trying OCR for {file_path}")
                ocr_text = FileService._extract_pdf_with_ocr(file_path)
                if ocr_text and len(ocr_text) > len(full_text or ""):
                    full_text = ocr_text
                    logger.info(f"OCR extracted {len(full_text)} chars from {file_path}")

            if not full_text or not full_text.strip():
                logger.warning(f"No text extracted from PDF: {file_path}")
                return None

            return full_text

        except Exception as e:
            logger.error(f"Error extracting text from PDF {file_path}: {str(e)}")
            return None

    @staticmethod
    def extract_docx_text(file_path: str) -> Optional[str]:
        """
        Extract text from a Word document (.docx).

        Args:
            file_path: Path to the DOCX file

        Returns:
            Extracted text as string, or None if extraction fails
        """
        if not DOCX_AVAILABLE:
            logger.warning("python-docx not installed, cannot extract DOCX text")
            return None

        try:
            doc_file = Path(file_path)
            if not doc_file.exists():
                logger.error(f"DOCX file not found: {file_path}")
                return None

            doc = DocxDocument(file_path)
            text_content = []

            # Extract text from paragraphs
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_content.append(paragraph.text)

            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        text_content.append(" | ".join(row_text))

            full_text = "\n".join(text_content)

            if not full_text.strip():
                logger.warning(f"No text extracted from DOCX: {file_path}")
                return None

            logger.info(f"Successfully extracted {len(full_text)} characters from DOCX {file_path}")
            return full_text

        except Exception as e:
            logger.error(f"Error extracting text from DOCX {file_path}: {str(e)}")
            return None

    @staticmethod
    def extract_pptx_text(file_path: str) -> Optional[str]:
        """
        Extract text from a PowerPoint presentation (.pptx).

        Args:
            file_path: Path to the PPTX file

        Returns:
            Extracted text as string, or None if extraction fails
        """
        if not PPTX_AVAILABLE:
            logger.warning("python-pptx not installed, cannot extract PPTX text")
            return None

        try:
            pptx_file = Path(file_path)
            if not pptx_file.exists():
                logger.error(f"PPTX file not found: {file_path}")
                return None

            prs = Presentation(file_path)
            text_content = []

            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)

                if slide_text:
                    text_content.append(f"--- Slide {slide_num} ---")
                    text_content.extend(slide_text)

            full_text = "\n".join(text_content)

            if not full_text.strip():
                logger.warning(f"No text extracted from PPTX: {file_path}")
                return None

            logger.info(f"Successfully extracted {len(full_text)} characters from PPTX {file_path}")
            return full_text

        except Exception as e:
            logger.error(f"Error extracting text from PPTX {file_path}: {str(e)}")
            return None

    @staticmethod
    def extract_xlsx_text(file_path: str) -> Optional[str]:
        """
        Extract text from an Excel file (.xlsx).

        Args:
            file_path: Path to the XLSX file

        Returns:
            Extracted text as string, or None if extraction fails
        """
        if not XLSX_AVAILABLE:
            logger.warning("openpyxl not installed, cannot extract XLSX text")
            return None

        try:
            xlsx_file = Path(file_path)
            if not xlsx_file.exists():
                logger.error(f"XLSX file not found: {file_path}")
                return None

            workbook = openpyxl.load_workbook(file_path, data_only=True)
            text_content = []

            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                sheet_text = [f"--- Sheet: {sheet_name} ---"]

                for row in sheet.iter_rows():
                    row_values = []
                    for cell in row:
                        if cell.value is not None:
                            row_values.append(str(cell.value))
                    if row_values:
                        sheet_text.append(" | ".join(row_values))

                if len(sheet_text) > 1:  # More than just the header
                    text_content.extend(sheet_text)

            full_text = "\n".join(text_content)

            if not full_text.strip():
                logger.warning(f"No text extracted from XLSX: {file_path}")
                return None

            logger.info(f"Successfully extracted {len(full_text)} characters from XLSX {file_path}")
            return full_text

        except Exception as e:
            logger.error(f"Error extracting text from XLSX {file_path}: {str(e)}")
            return None

    @staticmethod
    def get_pptx_slide_count(file_path: str) -> Optional[int]:
        """
        Get the number of slides in a PowerPoint presentation.

        Args:
            file_path: Path to the PPTX file

        Returns:
            Number of slides, or None if reading fails
        """
        if not PPTX_AVAILABLE:
            return None

        try:
            pptx_file = Path(file_path)
            if not pptx_file.exists():
                return None

            prs = Presentation(file_path)
            return len(prs.slides)

        except Exception as e:
            logger.error(f"Error getting slide count from PPTX {file_path}: {str(e)}")
            return None

    @staticmethod
    def get_docx_page_count(file_path: str) -> Optional[int]:
        """
        Estimate page count for a Word document.
        Note: This is an estimate based on paragraph count since DOCX doesn't store page count.

        Args:
            file_path: Path to the DOCX file

        Returns:
            Estimated page count, or None if reading fails
        """
        if not DOCX_AVAILABLE:
            return None

        try:
            doc_file = Path(file_path)
            if not doc_file.exists():
                return None

            doc = DocxDocument(file_path)
            # Rough estimate: ~25 paragraphs per page
            paragraph_count = len([p for p in doc.paragraphs if p.text.strip()])
            return max(1, paragraph_count // 25 + 1)

        except Exception as e:
            logger.error(f"Error estimating page count from DOCX {file_path}: {str(e)}")
            return None

    @staticmethod
    def extract_file_text(file_path: str, file_extension: str) -> Optional[str]:
        """
        Extract text from a file based on its extension.

        Args:
            file_path: Path to the file
            file_extension: File extension (e.g., '.pdf', '.docx')

        Returns:
            Extracted text as string, or None if extraction fails or not supported
        """
        # Normalize extension to lowercase
        ext = file_extension.lower()

        if ext == ".pdf":
            text = FileService.extract_pdf_text(file_path)
        elif ext in [".docx", ".doc"]:
            text = FileService.extract_docx_text(file_path)
        elif ext in [".pptx", ".ppt"]:
            text = FileService.extract_pptx_text(file_path)
        elif ext in [".xlsx", ".xls"]:
            text = FileService.extract_xlsx_text(file_path)
        elif ext == ".txt":
            try:
                with open(file_path, 'r', encoding='utf-8') as f:
                    text = f.read()
            except Exception as e:
                logger.error(f"Error reading TXT file {file_path}: {str(e)}")
                return None
        else:
            logger.info(f"Text extraction not supported for file type: {ext}")
            return None

        # PostgreSQL text fields cannot contain NUL (0x00) bytes
        return text.replace('\x00', '') if text else text

    @staticmethod
    def get_file_page_count(file_path: str, file_extension: str) -> Optional[int]:
        """
        Get page/slide count for a file based on its extension.

        Args:
            file_path: Path to the file
            file_extension: File extension

        Returns:
            Page/slide count, or None if not applicable or fails
        """
        ext = file_extension.lower()

        if ext == ".pdf":
            return FileService.get_pdf_page_count(file_path)
        elif ext in [".docx", ".doc"]:
            return FileService.get_docx_page_count(file_path)
        elif ext in [".pptx", ".ppt"]:
            return FileService.get_pptx_slide_count(file_path)

        return None
